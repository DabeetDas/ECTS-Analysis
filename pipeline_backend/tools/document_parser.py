from __future__ import annotations

from pathlib import Path

from pipeline_backend.schemas import TranscriptChunk


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in {".txt", ".md", ".csv", ".json"}:
        return path.read_text(encoding="utf-8", errors="ignore")
    if suffix == ".pdf":
        return extract_pdf_text(path)
    if suffix == ".pptx":
        return extract_pptx_text(path)
    raise ValueError(f"Unsupported document type: {path.suffix}")


def extract_pdf_text(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise ImportError("PDF parsing requires `pypdf`. Add it to the pipeline backend environment.") from exc
    reader = PdfReader(str(path))
    return "\n\n".join(page.extract_text() or "" for page in reader.pages)


def extract_pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ImportError("PPTX parsing requires `python-pptx`. Add it to the pipeline backend environment.") from exc
    presentation = Presentation(str(path))
    chunks: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                texts.append(text)
        if texts:
            chunks.append(f"[Slide {slide_index}]\n" + "\n".join(texts))
    return "\n\n".join(chunks)


def chunk_transcript(text: str, *, source_name: str | None = None, chunk_chars: int = 3500) -> list[TranscriptChunk]:
    paragraphs = [paragraph.strip() for paragraph in text.splitlines() if paragraph.strip()]
    chunks: list[TranscriptChunk] = []
    buffer: list[str] = []
    size = 0
    for paragraph in paragraphs:
        if buffer and size + len(paragraph) > chunk_chars:
            chunks.append(
                TranscriptChunk(
                    chunk_id=f"chunk_{len(chunks) + 1}",
                    text="\n".join(buffer),
                    source_name=source_name,
                )
            )
            buffer = []
            size = 0
        buffer.append(paragraph)
        size += len(paragraph) + 1
    if buffer:
        chunks.append(
            TranscriptChunk(
                chunk_id=f"chunk_{len(chunks) + 1}",
                text="\n".join(buffer),
                source_name=source_name,
            )
        )
    return chunks

